#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
生成自媒体IP定位介绍PPT
作者：林琨奇
日期：2026-05-08
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# 创建空白演示文稿
prs = Presentation()

# ========== 第1页：封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

# 设置背景颜色（深蓝色）
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(26, 42, 68)

# 标题
title_box = slide.shapes.title
if not title_box:
    title_box = slide.shapes.placeholders[0] if slide.shapes.placeholders else None

title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
tf = title.text_frame
p = tf.add_paragraph()
p.text = "我是林kunki"
p.font.size = Pt(54)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
p.font.bold = True

# 副标题
subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(8), Inches(0.8))
tf = subtitle.text_frame
p = tf.add_paragraph()
p.text = "AI时代一人公司实践者"
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(200, 220, 255)
p.alignment = PP_ALIGN.CENTER

# 标注
footer = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.3))
tf = footer.text_frame
p = tf.add_paragraph()
p.text = "540学生法则 · 边保研边创业"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(150, 170, 200)
p.alignment = PP_ALIGN.CENTER

# ========== 第2页：关于我 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

# 标题
title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
tf = title.text_frame
p = tf.add_paragraph()
p.text = "关于我"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(26, 42, 68)

# 内容
content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5))
tf = content.text_frame

points = [
    "▫️ 05后 | 广外国贸大三 | 专业第一",
    "▫️ SCI二区论文一作 | 亚太AI大赛金奖",
    "▫️ 一边准备保研，一边搭建AI一人公司",
    "▫️ 和北京投资人合作写《OpenClaw智能体运营》",
    "",
    "顶层定位：**AI时代一人公司实践者**",
    "落地体系：**540学生法则**"
]

for point in points:
    p = tf.add_paragraph()
    p.text = point
    if "**" in point:
        p.font.bold = True
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(45, 65, 102)
    else:
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
    p.space_after = Pt(12)

# ========== 第3页：核心哲学：540学生法则 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
tf = title.text_frame
p = tf.add_paragraph()
p.text = "核心哲学：540学生法则"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(26, 42, 68)

content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5))
tf = content.text_frame

p = tf.add_paragraph()
p.text = "540 = 三个100 + 四个60"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RGBColor(45, 65, 102)
p.space_after = Pt(20)

p = tf.add_paragraph()
p.text = "三个100（底层系统构建）"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(80, 80, 80)
p.space_after = Pt(8)

points_3 = [
    "✅ 100小时AI工具使用 —— 真刀真枪解决问题",
    "✅ 100小时目标试错推进 —— 集中打穿拿到结果",
    "✅ 100小时优质信息输入 —— 停止输入就是停止进化"
]
for point in points_3:
    p = tf.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.level = 1
    p.space_after = Pt(6)

p = tf.add_paragraph()
p.text = "\n四个60（现实能力训练场）"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = RGBColor(80, 80, 80)
p.space_after = Pt(8)

points_4 = [
    "💪 60天健身 —— 练自律、精力、抗压",
    "🗣️ 60天社交表达 —— 链接、破冰、价值传递",
    "📹 60天内容输出 —— 练表达、建数字资产",
    "💰 60天销售成交 —— 练需求洞察、信任、闭环"
]
for point in points_4:
    p = tf.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.level = 1
    p.space_after = Pt(6)

p = tf.add_paragraph()
p.text = "\n一句话：课内用AI高效拿结果，课外做540练能力"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(45, 65, 102)
p.font.italic = True

# ========== 第4页：人设三身份 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
tf = title.text_frame
p = tf.add_paragraph()
p.text = "人设三身份"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(26, 42, 68)

content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5))
tf = content.text_frame

identities = [
    ("1. 真实实践者", "不端着，持续记录一步步打顺路", RGBColor(60, 120, 200)),
    ("2. 结果拆解者", "不只讲故事，把方法全拆出来给你看", RGBColor(60, 180, 120)),
    ("3. 路径参考者", "不替你做决定，给可复制可调整的路径", RGBColor(180, 90, 160)),
]

for name, desc, color in identities:
    p = tf.add_paragraph()
    p.text = name
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = color
    p.space_after = Pt(4)

    p2 = tf.add_paragraph()
    p2.text = "   " + desc
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(60, 60, 60)
    p2.space_after = Pt(20)

# ========== 第5页：差异化优势 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
tf = title.text_frame
p = tf.add_paragraph()
p.text = "差异化优势"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(26, 42, 68)

content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5))
tf = content.text_frame

diff_points = [
    ("05后大学生身份", "年纪最小结果最硬，2-3年窗口期",),
    ("校内外双线硬结果", "绩点第一+SCI+竞赛 + AI商业化从0到1，组合稀缺",),
    ("正在发生的连续叙事", "从象牙塔走向市场，别人抄不走",),
    ("540法则专属方法论", "所有内容产品都挂在这个标签下",),
]

for i, (point, desc) in enumerate(diff_points):
    p = tf.add_paragraph()
    p.text = f"{i+1}. {point}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(45, 65, 102)
    p.space_after = Pt(3)

    p2 = tf.add_paragraph()
    p2.text = f"   {desc}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(60, 60, 60)
    p2.space_after = Pt(16)

# ========== 第6页：内容四条主线 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
tf = title.text_frame
p = tf.add_paragraph()
p.text = "内容四条主线"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(26, 42, 68)

content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5))
tf = content.text_frame

content_lines = [
    ("大学结果", "绩点、论文、竞赛、保研、学习系统\n核心：讲如何用最低成本拿到关键结果"),
    ("AI工具与系统", "真实用过的AI工具、智能体、工作流、OpenClaw\n核心：只讲怎么用、适合谁、解决什么问题"),
    ("一人公司实践", "内容输出、个人网站、私域、产品、品牌搭建\n核心：全过程真实记录"),
    ("成长实验", "健身、社交、行动力、形象管理\n核心：不写鸡汤，写实验记录和方法复盘"),
]

emojis = ["📚", "🤖", "🏭", "🧪"]
for i, ((title_text, core), emoji) in enumerate(zip(content_lines, emojis)):
    p = tf.add_paragraph()
    p.text = f"{emoji} {title_text}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(45, 65, 102)
    p.space_after = Pt(4)

    p2 = tf.add_paragraph()
    p2.text = f"   {core}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(60, 60, 60)
    p2.space_after = Pt(14)

# ========== 第7页：商业模式：四层漏斗 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
tf = title.text_frame
p = tf.add_paragraph()
p.text = "商业模式：四层漏斗"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(26, 42, 68)

content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5))
tf = content.text_frame

funnel = [
    ("🎁 免费层", "公域内容+网站+540思维导图+AI工具清单", "免费", "获客，建立第一道信任"),
    ("💸 低价层", "《540入门工具包》", "¥9.9", "筛选付费意愿，沉淀种子用户"),
    ("👥 中价层", "《540年度陪伴社群》", "¥199/年", "主营收，每周分享践行方法"),
    ("🚀 高价层", "《一人公司起步陪跑》", "¥1299", "1v1帮定位、搭工作流、找方向"),
]

for name, product, price, role in funnel:
    p = tf.add_paragraph()
    p.text = f"{name} | {product}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(45, 65, 102)
    p.space_after = Pt(2)

    p2 = tf.add_paragraph()
    p2.text = f"   价格：{price} | 作用：{role}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(60, 60, 60)
    p2.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "\n一句话：用内容获客，用网站建信任，用代充做第一次低门槛成交，用自有产品做主营收"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(45, 65, 102)
p.font.italic = True

# ========== 第8页：AI智能体分工：六人团队 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
tf = title.text_frame
p = tf.add_paragraph()
p.text = "AI分工：六人自动运转团队"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(26, 42, 68)

content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5))
tf = content.text_frame

team = [
    ("选哥", "每日选题推荐官", "爬热点→推选题→评分→写入飞书"),
    ("盘哥", "复盘沉淀官", "录数据→写复盘→沉淀素材→更新记忆"),
    ("催哥", "进度催促官", "每日提醒任务，检查进度打气"),
    ("情报哥", "信息收集官", "不定期收集院校/行业情报"),
    ("写作哥", "内容写作官", "根据选题写出完整初稿"),
    ("Main Agent", "系统运维枢纽", "飞书同步+全链路创作调度"),
]

for name, role, desc in team:
    p = tf.add_paragraph()
    p.text = f"• {name} - {role}"
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = RGBColor(45, 65, 102)
    p.space_after = Pt(2)

    p2 = tf.add_paragraph()
    p2.text = f"  {desc}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(60, 60, 60)
    p2.space_after = Pt(10)

# ========== 第9页：持续迭代飞轮 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(250, 250, 250)

title = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(0.5))
tf = title.text_frame
p = tf.add_paragraph()
p.text = "持续迭代：自动飞轮闭环"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(26, 42, 68)

content = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5))
tf = content.text_frame

steps = [
    "1️⃣ 飞书选哥：每日自动爬热点，推荐TOP3选题",
    "2️⃣ 林总选择选题，AI协助13步生产出文案",
    "3️⃣ 发布到平台，获取播放/点赞/评论数据",
    "4️⃣ 飞书盘哥：录入数据，复盘成败，沉淀素材",
    "5️⃣ 沉淀经验到知识库，更新选题命中率",
    "6️⃣ 次日选哥推荐更精准，系统越滚越聪明"
]

p = tf.add_paragraph()
p.text = "这是一个活的系统，越用越聪明："
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(60, 60, 60)
p.space_after = Pt(15)

for step in steps:
    p = tf.add_paragraph()
    p.text = step
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "\n核心优势：AI代替99%重复劳动，1%创造力永远属于你"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(45, 65, 102)
p.font.bold = True

# ========== 第10页：结尾页 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(26, 42, 68)

title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
tf = title.text_frame
p = tf.add_paragraph()
p.text = "感谢聆听"
p.font.size = Pt(48)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
p.font.bold = True

subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
tf = subtitle.text_frame
p = tf.add_paragraph()
p.text = "540学生法则 · AI一人公司实践者"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(200, 220, 255)
p.alignment = PP_ALIGN.CENTER

# 保存PPT
output_path = "D:/ObsidianVaults/MyVault/E-产出交付/06-临时文件/我是林kunki-自媒体IP定位介绍.pptx"
prs.save(output_path)
print(f"PPT已生成：{output_path}")
